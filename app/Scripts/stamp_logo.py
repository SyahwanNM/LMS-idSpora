import os
import sys
import argparse
from PIL import Image

def stamp_pptx(pptx_path, logo_path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Error: python-pptx is not installed.")
        sys.exit(1)
        
    prs = Presentation(pptx_path)
    
    # Get aspect ratio of logo
    logo_img = Image.open(logo_path)
    ratio = logo_img.height / logo_img.width
    
    for slide in prs.slides:
        # Standard width of logo: 1.5 inches
        width = Inches(1.5)
        height = Inches(1.5 * ratio)
        
        # Top-right corner
        left = prs.slide_width - width - Inches(0.4)
        top = Inches(0.4)
        
        slide.shapes.add_picture(logo_path, left, top, width=width, height=height)
        
    prs.save(pptx_path)
    print("PPTX stamped successfully.")

def stamp_pdf(pdf_path, logo_path):
    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
    except ImportError:
        print("Error: pypdf or reportlab is not installed.")
        sys.exit(1)
        
    import tempfile
    
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    
    # Open logo image to get aspect ratio
    logo_img = Image.open(logo_path)
    ratio = logo_img.height / logo_img.width
    
    for i in range(len(reader.pages)):
        page = reader.pages[i]
        page_width = float(page.mediabox.width)
        page_height = float(page.mediabox.height)
        
        fd, temp_pdf = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        
        try:
            # We want logo width to be 15% of the page width
            logo_w = page_width * 0.15
            logo_h = logo_w * ratio
            
            c = canvas.Canvas(temp_pdf, pagesize=(page_width, page_height))
            # Margins (3% of page width)
            margin = page_width * 0.03
            x = page_width - logo_w - margin
            y = page_height - logo_h - margin
            
            c.drawImage(logo_path, x, y, width=logo_w, height=logo_h, mask='auto')
            c.save()
            
            watermark_reader = PdfReader(temp_pdf)
            watermark_page = watermark_reader.pages[0]
            
            page.merge_page(watermark_page, over=True)
            writer.add_page(page)
        except Exception as e:
            print(f"Error watermarking page {i}: {e}")
            writer.add_page(page)
        finally:
            if os.path.exists(temp_pdf):
                os.remove(temp_pdf)
                
    with open(pdf_path, "wb") as f:
        writer.write(f)
    print("PDF stamped successfully.")

def stamp_image(img_path, logo_path):
    img = Image.open(img_path)
    logo = Image.open(logo_path)
    
    # Resize logo to 15% of image width
    logo_w = int(img.width * 0.15)
    logo_h = int(logo_w * logo.height / logo.width)
    logo = logo.resize((logo_w, logo_h), Image.Resampling.LANCZOS)
    
    # Coordinates for top-right corner
    margin_w = int(img.width * 0.03)
    margin_h = int(img.height * 0.03)
    x = img.width - logo_w - margin_w
    y = margin_h
    
    # Paste logo
    if logo.mode == 'RGBA':
        # Paste with transparent mask
        img.paste(logo, (x, y), logo)
    else:
        img.paste(logo, (x, y))
        
    img.save(img_path)
    print("Image stamped successfully.")

def main():
    parser = argparse.ArgumentParser(description="Stamp idSpora logo on slides/documents/images")
    parser.add_argument("--file", required=True, help="Path to target file")
    parser.add_argument("--logo", required=True, help="Path to logo file")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.file):
        print(f"Error: File '{args.file}' not found.")
        sys.exit(1)
        
    if not os.path.exists(args.logo):
        print(f"Error: Logo '{args.logo}' not found.")
        sys.exit(1)
        
    ext = os.path.splitext(args.file)[1].lower()
    
    if ext == '.pptx':
        stamp_pptx(args.file, args.logo)
    elif ext == '.pdf':
        stamp_pdf(args.file, args.logo)
    elif ext in ['.png', '.jpg', '.jpeg']:
        stamp_image(args.file, args.logo)
    else:
        print(f"Error: Unsupported file format '{ext}'")
        sys.exit(1)

if __name__ == "__main__":
    main()
